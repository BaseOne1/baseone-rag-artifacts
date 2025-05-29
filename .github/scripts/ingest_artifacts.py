import os
import uuid
import boto3
import pdfplumber
import pandas as pd
import requests
from docx import Document
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
from decimal import Decimal
import hashlib


dynamodb = boto3.resource("dynamodb", region_name=os.environ["AWS_REGION"])
table = dynamodb.Table("BaseOneRAG")
OPENAI_API_KEY = os.environ["OPENAI_API_KEY"]

text_splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)

def embed_text(text):
    try:
        response = requests.post(
            "https://api.openai.com/v1/embeddings",
            headers={"Authorization": f"Bearer {OPENAI_API_KEY}", "Content-Type": "application/json"},
            json={"input": text, "model": "text-embedding-3-small"}
        )
        response.raise_for_status()
        data = response.json()
        return data["data"][0]["embedding"]
    except Exception as e:
        print(f"[ERROR] Failed to embed text: {e}")
        return None
    
def compute_file_hash(filepath):
    hasher = hashlib.sha256()
    with open(filepath, 'rb') as f:
        buf = f.read()
        hasher.update(buf)
    return hasher.hexdigest()

def extract_text_from_file(filepath):
    ext = filepath.lower().split(".")[-1]
    text = ""
    try:
        if ext == "pdf":
            with pdfplumber.open(filepath) as pdf:
                text = "\n".join(page.extract_text() or "" for page in pdf.pages)
        elif ext in ["xls", "xlsx"]:
            df = pd.read_excel(filepath, engine="openpyxl")
            text = df.astype.str().apply(lambda x: " | ".join(x), axis=1).str.cat(sep="\n")
        elif ext == "docx":
            doc = Document(filepath)
            text = "\n".join([para.text for para in doc.paragraphs])
        elif ext == "pptx":
            ppt = Presentation(filepath)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif ext == "txt":
            with open(filepath, "r", encoding="utf-8") as f:
                text = f.read()
    except Exception as e:
        print(f"[ERROR] Failed to extract {filepath}: {e}")
    return text

def chunk_and_upload(text, source, file_hash):
    if not text.strip():
        return

    # Check for existing file by hash
    existing = table.query(
        KeyConditionExpression=boto3.dynamodb.conditions.Key('PK').eq('DOC#BaseOne')
    ).get('Items', [])

    if any(item.get('file_hash') == file_hash for item in existing):
        print(f"[SKIP] {source} already ingested.")
        return

    # Process and upload
    chunks = text_splitter.split_text(text)
    for chunk in chunks:
        embedding = embed_text(chunk)
        if embedding:
            embedding_decimal = [Decimal(str(x)) for x in embedding]
            id = str(uuid.uuid4())
            try:
                table.put_item(Item={
                    "PK": "DOC#BaseOne",
                    "SK": f"CHUNK#{id}",
                    "text": chunk,
                    "embedding": embedding_decimal,
                    "source": source,
                    "file_hash": file_hash
                })
                print(f"Uploaded chunk from {source} with ID {id}")
            except Exception as e:
                print(f"[ERROR] Failed to upload chunk from {source}: {e}")


def process_all_files(folder=".github/artifacts"):
    print(f"[INFO] Scanning folder: {folder}")
    for root, _, files in os.walk(folder):
        for file in files:
            path = os.path.join(root, file)
            print(f"Processing {path}...")
            file_hash = compute_file_hash(path)
            text = extract_text_from_file(path)
            chunk_and_upload(text, source=path, file_hash=file_hash)


if __name__ == "__main__":
    process_all_files()
